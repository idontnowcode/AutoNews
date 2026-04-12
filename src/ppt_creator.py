from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# 9:16 세로형 (1080 x 1920 기준)
SLIDE_W = Inches(6.75)   # 6.75 / 12.0 = 9/16
SLIDE_H = Inches(12.0)

COLORS = {
    'bg_title': RGBColor(0x1F, 0x4E, 0x79),   # 진한 파랑
    'bg_body':  RGBColor(0xF7, 0xF9, 0xFF),   # 연한 흰색
    'white':    RGBColor(0xFF, 0xFF, 0xFF),
    'dark':     RGBColor(0x1A, 0x1A, 0x2E),
    'accent':   RGBColor(0x2E, 0x75, 0xB6),
}


def _add_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, text, left, top, width, height,
                 font_size, bold=False, color=None, align=PP_ALIGN.CENTER):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color or COLORS['dark']


def _make_title_slide(prs, title: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, COLORS['bg_title'])

    # 상단 장식바
    bar = slide.shapes.add_shape(1, Inches(0), Inches(4.5), SLIDE_W, Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0x2E, 0xCC, 0xCC)
    bar.line.fill.background()

    # 제목 텍스트
    _add_textbox(slide, title, Inches(0.4), Inches(5.0),
                 SLIDE_W - Inches(0.8), Inches(2.5),
                 font_size=Pt(36), bold=True, color=COLORS['white'])

    # 하단 구분선
    _add_textbox(slide, 'Shorts  |  오늘의 뉴스', Inches(0.4), Inches(10.5),
                 SLIDE_W - Inches(0.8), Inches(0.8),
                 font_size=Pt(18), color=RGBColor(0xAA, 0xCC, 0xFF))


def _make_content_slide(prs, slide_data: dict, slide_num: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, COLORS['bg_body'])

    # 슬라이드 번호 뱃지
    badge = slide.shapes.add_shape(1, Inches(0.3), Inches(0.5), Inches(0.5), Inches(0.5))
    badge.fill.solid()
    badge.fill.fore_color.rgb = COLORS['accent']
    badge.line.fill.background()

    # 이모지 + 제목
    heading = f"{slide_data.get('emoji', '')} {slide_data['heading']}"
    _add_textbox(slide, heading, Inches(0.3), Inches(1.5),
                 SLIDE_W - Inches(0.6), Inches(2.0),
                 font_size=Pt(34), bold=True, color=COLORS['accent'])

    # 본문
    _add_textbox(slide, slide_data['body'], Inches(0.3), Inches(4.0),
                 SLIDE_W - Inches(0.6), Inches(5.0),
                 font_size=Pt(28), align=PP_ALIGN.LEFT)


def create_ppt(script_data: dict, output_path: str) -> str:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    _make_title_slide(prs, script_data['title'])
    for i, s in enumerate(script_data.get('slides', [])[:3], 1):
        _make_content_slide(prs, s, i)

    prs.save(output_path)
    return output_path
